#!/usr/bin/env python3
"""Generate cost analysis PPT for sandbox runtime pod density benchmarks."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
import os

# ── AWS Pricing (us-west-2, On-Demand, Linux) ──
PRICES = {
    "m8i.2xlarge": 0.4234,
    "m7g.2xlarge": 0.3264,
    "m7g.metal": 2.6112,
    "c7g.metal": 2.3200,
}
HOURS_PER_MONTH = 730

# ── Colors ──
DARK_BG = RGBColor(0x1B, 0x1B, 0x2F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x00)
ACCENT_GREEN = RGBColor(0x00, 0xCC, 0x88)
ACCENT_BLUE = RGBColor(0x44, 0x99, 0xFF)
ACCENT_RED = RGBColor(0xFF, 0x55, 0x55)
ACCENT_PURPLE = RGBColor(0xBB, 0x77, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xDD, 0xDD)
ACCENT_PINK = RGBColor(0xFF, 0x77, 0xAA)
DIM_GRAY = RGBColor(0x88, 0x88, 0x99)

BAR_COLORS = [ACCENT_RED, ACCENT_ORANGE, ACCENT_BLUE, ACCENT_CYAN, ACCENT_GREEN, ACCENT_PURPLE, ACCENT_PINK]

# ── Test Data ──
scenarios = [
    {"short": "m8i.2x\nkata-qemu",   "instance": "m8i.2xlarge", "runtime": "kata-qemu",  "pods": 7,  "theoretical": 12, "bottleneck": "VM crash (restart)", "arch": "Intel x86", "vcpu": 8, "mem": 32, "overhead": "100m/250Mi", "isolation": "嵌套 VM", "utilization": "58%"},
    {"short": "m8i.2x\nkata-clh",    "instance": "m8i.2xlarge", "runtime": "kata-clh",   "pods": 13, "theoretical": 13, "bottleneck": "调度器内存上限",    "arch": "Intel x86", "vcpu": 8, "mem": 32, "overhead": "100m/200Mi", "isolation": "嵌套 VM", "utilization": "100%"},
    {"short": "m8i.2x\ngVisor",      "instance": "m8i.2xlarge", "runtime": "gVisor",     "pods": 14, "theoretical": 14, "bottleneck": "调度器内存上限",    "arch": "Intel x86", "vcpu": 8, "mem": 32, "overhead": "无",         "isolation": "用户态内核", "utilization": "100%"},
    {"short": "m7g.2x\ngVisor+EFS",  "instance": "m7g.2xlarge", "runtime": "gVisor",     "pods": 14, "theoretical": 14, "bottleneck": "调度器内存上限",    "arch": "Graviton",  "vcpu": 8, "mem": 32, "overhead": "无",         "isolation": "用户态内核", "utilization": "100%"},
    {"short": "m7g.metal\nkata-fc",   "instance": "m7g.metal",   "runtime": "kata-fc",    "pods": 90, "theoretical": 91, "bottleneck": "CPU 99%",          "arch": "Graviton",  "vcpu": 64, "mem": 256, "overhead": "250m/130Mi", "isolation": "裸金属 FC", "utilization": "99%"},
    {"short": "c7g.metal\nkata-fc",   "instance": "c7g.metal",   "runtime": "kata-fc",    "pods": 55, "theoretical": 55, "bottleneck": "内存 99%",          "arch": "Graviton",  "vcpu": 64, "mem": 128, "overhead": "250m/130Mi", "isolation": "裸金属 FC", "utilization": "100%"},
    {"short": "m7g.metal\nFC+EBS PVC","instance": "m7g.metal",   "runtime": "kata-fc",    "pods": 28, "theoretical": 28, "bottleneck": "EBS 挂载上限 (31 shared)", "arch": "Graviton", "vcpu": 64, "mem": 256, "overhead": "250m/130Mi+PVC", "isolation": "裸金属 FC", "utilization": "100%"},
]

for s in scenarios:
    h = PRICES[s["instance"]]
    s["hourly"] = h
    s["monthly"] = h * HOURS_PER_MONTH
    s["per_pod_hr"] = h / s["pods"]
    s["per_pod_mo"] = h * HOURS_PER_MONTH / s["pods"]

sorted_by_cost = sorted(scenarios, key=lambda x: x["per_pod_hr"])

# ── Helpers ──
def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    return table_shape.table

def style_cell(cell, text, font_size=10, bold=False, color=WHITE, bg_color=None, alignment=PP_ALIGN.CENTER):
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color

# ── Create Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

add_text_box(slide, 1.5, 1.5, 10, 1.2,
    "Sandbox Runtime Cost Analysis", 36, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 2.8, 10, 0.8,
    "Per-Pod Cost Comparison: kata-qemu vs kata-clh vs gVisor vs kata-fc (Firecracker)", 18, False, LIGHT_GRAY, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 4.0, 10, 0.5,
    "AWS EKS · us-west-2 · On-Demand Pricing · April 2026", 14, False, DIM_GRAY, PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.0, 10, 1.5,
    "Pod Spec: 4 containers × stress-ng (gateway 150m/1G + config-watcher 100m/256M + envoy 100m/256M + wazuh 100m/512M)\n"
    "Total per pod: 450m CPU / 2048 MiB memory (Guaranteed QoS, request = limit)\n"
    "All scenarios: full CPU + memory stress, zero idle pods",
    12, False, DIM_GRAY, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# SLIDE 2: Cost Per Pod Bar Chart (horizontal bars using shapes)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.2, 12, 0.6,
    "Per-Pod Monthly Cost — Lower is Better", 28, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 0.85, 12, 0.4,
    "Instance cost ÷ max stable pods = cost per sandbox instance per month", 13, False, DIM_GRAY, PP_ALIGN.CENTER)

# Draw horizontal bar chart
bar_area_left = 3.2
bar_area_top = 1.5
bar_area_width = 8.5
bar_height = 0.55
bar_gap = 0.18
max_cost = max(s["per_pod_mo"] for s in sorted_by_cost) * 1.15

for i, s in enumerate(sorted_by_cost):
    y = bar_area_top + i * (bar_height + bar_gap)
    
    # Label on left
    label = s["short"].replace("\n", " ")
    add_text_box(slide, 0.2, y - 0.05, 2.9, bar_height + 0.1,
        label, 11, True, WHITE, PP_ALIGN.RIGHT)
    
    # Bar
    bar_width = (s["per_pod_mo"] / max_cost) * bar_area_width
    shape = slide.shapes.add_shape(
        1,  # rectangle
        Inches(bar_area_left), Inches(y),
        Inches(bar_width), Inches(bar_height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BAR_COLORS[i % len(BAR_COLORS)]
    shape.line.fill.background()
    
    # Value label
    add_text_box(slide, bar_area_left + bar_width + 0.1, y - 0.02, 2, bar_height,
        f"${s['per_pod_mo']:.2f}/mo", 12, True, BAR_COLORS[i % len(BAR_COLORS)], PP_ALIGN.LEFT)
    
    # Pod count + instance
    add_text_box(slide, bar_area_left + bar_width + 0.1, y + 0.25, 2.5, 0.3,
        f"{s['pods']} pods · {s['instance']}", 9, False, DIM_GRAY, PP_ALIGN.LEFT)

# Notes
add_text_box(slide, 0.5, 6.7, 12, 0.6,
    "Note: On-Demand pricing, us-west-2. RI/Savings Plans can reduce by 30-60%. Graviton instances are ~23% cheaper than Intel equivalents.",
    9, False, DIM_GRAY, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════
# SLIDE 3: Detailed Comparison Table
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.3, 0.15, 12, 0.6,
    "Full Comparison Matrix", 28, True, WHITE, PP_ALIGN.CENTER)

# Table
cols = 10
rows = len(scenarios) + 1
tbl = add_table(slide, 0.3, 0.85, 12.7, 5.5, rows, cols)

# Headers
headers = ["场景", "实例", "架构", "vCPU/RAM", "运行时", "Max Pods", "理论值", "瓶颈", "$/Pod/hr", "$/Pod/月"]
header_bg = RGBColor(0x2A, 0x2A, 0x45)
for j, h in enumerate(headers):
    style_cell(tbl.cell(0, j), h, 10, True, ACCENT_ORANGE, header_bg)

# Data rows
row_bg1 = RGBColor(0x22, 0x22, 0x38)
row_bg2 = RGBColor(0x28, 0x28, 0x42)
for i, s in enumerate(scenarios):
    bg = row_bg1 if i % 2 == 0 else row_bg2
    row = i + 1
    style_cell(tbl.cell(row, 0), s["short"].replace("\n", " "), 9, True, WHITE, bg)
    style_cell(tbl.cell(row, 1), s["instance"], 9, False, LIGHT_GRAY, bg)
    style_cell(tbl.cell(row, 2), s["arch"], 9, False, LIGHT_GRAY, bg)
    style_cell(tbl.cell(row, 3), f"{s['vcpu']}v / {s['mem']}G", 9, False, LIGHT_GRAY, bg)
    style_cell(tbl.cell(row, 4), s["runtime"], 9, False, LIGHT_GRAY, bg)
    
    # Highlight best pods
    pod_color = ACCENT_GREEN if s["pods"] >= 55 else (ACCENT_BLUE if s["pods"] >= 13 else ACCENT_RED)
    style_cell(tbl.cell(row, 5), str(s["pods"]), 11, True, pod_color, bg)
    style_cell(tbl.cell(row, 6), str(s["theoretical"]), 9, False, DIM_GRAY, bg)
    style_cell(tbl.cell(row, 7), s["bottleneck"], 8, False, LIGHT_GRAY, bg)
    
    # Cost - highlight cheapest
    cost_color = ACCENT_GREEN if s["per_pod_hr"] < 0.030 else (WHITE if s["per_pod_hr"] < 0.050 else ACCENT_RED)
    style_cell(tbl.cell(row, 8), f"${s['per_pod_hr']:.4f}", 10, True, cost_color, bg)
    style_cell(tbl.cell(row, 9), f"${s['per_pod_mo']:.2f}", 10, True, cost_color, bg)

# Set column widths
widths = [1.8, 1.2, 0.9, 1.0, 1.0, 0.8, 0.7, 2.2, 1.0, 1.0]
for j, w in enumerate(widths):
    tbl.columns[j].width = Inches(w)

# Notes box
notes_text = (
    "测试参数: 每 Pod 4 容器 (gateway 150m/1G + config-watcher 100m/256M + envoy 100m/256M + wazuh 100m/512M), "
    "Guaranteed QoS, stress-ng full load (CPU 95% + vm-keep)\n"
    "Pod Overhead: kata-qemu 100m/250Mi · kata-clh 100m/200Mi · kata-fc 250m/130Mi · gVisor 无\n"
    "定价: AWS On-Demand, us-west-2, 730 hrs/month"
)
add_text_box(slide, 0.3, 6.5, 12.5, 0.9, notes_text, 9, False, DIM_GRAY, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════
# SLIDE 4: Same Instance Comparison (m8i.2xlarge)
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.2, 12, 0.6,
    "Same Instance, Different Runtimes — m8i.2xlarge (8 vCPU, 32 GiB)", 24, True, WHITE, PP_ALIGN.CENTER)

m8i_scenarios = [s for s in scenarios if s["instance"] == "m8i.2xlarge"]

# Big numbers
for i, s in enumerate(m8i_scenarios):
    x = 1.0 + i * 4.0
    color = [ACCENT_RED, ACCENT_BLUE, ACCENT_GREEN][i]
    
    # Runtime name
    add_text_box(slide, x, 1.2, 3.5, 0.5, s["runtime"], 22, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 1.8, 3.5, 0.4, s["isolation"], 12, False, DIM_GRAY, PP_ALIGN.CENTER)
    
    # Pod count - big
    add_text_box(slide, x, 2.4, 3.5, 1.0, str(s["pods"]), 60, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 3.4, 3.5, 0.4, "stable pods", 14, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    
    # Cost
    add_text_box(slide, x, 4.2, 3.5, 0.6, f"${s['per_pod_mo']:.2f}", 28, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 4.8, 3.5, 0.4, "per pod / month", 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    
    # Utilization
    add_text_box(slide, x, 5.5, 3.5, 0.4, f"理论达成: {s['utilization']}", 12, False, DIM_GRAY, PP_ALIGN.CENTER)
    add_text_box(slide, x, 5.9, 3.5, 0.4, f"瓶颈: {s['bottleneck']}", 10, False, DIM_GRAY, PP_ALIGN.CENTER)

# Key insight
add_text_box(slide, 0.5, 6.5, 12, 0.5,
    "💡 同一台 m8i.2xlarge: kata-qemu 每 Pod $44/月 → kata-clh $24/月 (↓46%) → gVisor $22/月 (↓50%)",
    14, True, ACCENT_ORANGE, PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 7.0, 12, 0.3,
    "Note: gVisor 安全隔离等级低于 Kata VM 隔离, syscall 兼容性有限制, EFS I/O 不可用",
    9, False, DIM_GRAY, PP_ALIGN.CENTER)

# ════════════════════════════════════════════════
# SLIDE 5: Bare Metal Scaling — kata-fc
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.2, 12, 0.6,
    "Bare Metal Scaling: kata-fc (Firecracker) on Graviton", 24, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 0.85, 12, 0.4,
    "裸金属 = 无嵌套虚拟化开销, Firecracker microVM 直接运行在硬件上", 13, False, DIM_GRAY, PP_ALIGN.CENTER)

fc_scenarios = [s for s in scenarios if "kata-fc" in s["runtime"]]
colors_fc = [ACCENT_PURPLE, ACCENT_CYAN, ACCENT_ORANGE]

for i, s in enumerate(fc_scenarios):
    x = 0.5 + i * 4.2
    color = colors_fc[i]
    
    inst_label = s["instance"]
    if "EBS PVC" in s["short"]:
        inst_label += " + EBS PVC"
    
    add_text_box(slide, x, 1.4, 3.8, 0.5, inst_label, 18, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 1.9, 3.8, 0.4, f"{s['vcpu']} vCPU · {s['mem']} GiB · {s['arch']}", 11, False, DIM_GRAY, PP_ALIGN.CENTER)
    
    # Pod count
    add_text_box(slide, x, 2.5, 3.8, 1.0, str(s["pods"]), 60, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 3.5, 3.8, 0.4, "stable pods (0 restarts)", 12, False, LIGHT_GRAY, PP_ALIGN.CENTER)
    
    # Cost
    add_text_box(slide, x, 4.2, 3.8, 0.6, f"${s['per_pod_mo']:.2f}/pod/月", 22, True, color, PP_ALIGN.CENTER)
    add_text_box(slide, x, 4.8, 3.8, 0.4, f"实例 ${s['hourly']:.4f}/hr", 11, False, DIM_GRAY, PP_ALIGN.CENTER)
    
    # Bottleneck
    add_text_box(slide, x, 5.4, 3.8, 0.4, f"瓶颈: {s['bottleneck']}", 11, False, DIM_GRAY, PP_ALIGN.CENTER)

# Key insight
add_text_box(slide, 0.5, 6.2, 12, 0.5,
    "💡 m7g.metal 90 pods = $21/pod/月 — 比嵌套虚拟化 kata-qemu 便宜 52%", 14, True, ACCENT_GREEN, PP_ALIGN.CENTER)
add_text_box(slide, 0.5, 6.7, 12, 0.5,
    "⚠️ 挂 EBS PVC 后仅 28 pods (EBS 共享槽位限制 31), 成本升至 $68/pod/月 — 比无 PVC 贵 3.2x", 13, True, ACCENT_RED, PP_ALIGN.CENTER)

add_text_box(slide, 0.5, 7.1, 12, 0.3,
    "Note: kata-fc 需要 devmapper thinpool, busybox 单层镜像限制, 不支持 CPU hotplug (aarch64)",
    9, False, DIM_GRAY, PP_ALIGN.LEFT)

# ════════════════════════════════════════════════
# SLIDE 6: Key Takeaways
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_text_box(slide, 0.5, 0.3, 12, 0.6,
    "Key Takeaways", 28, True, WHITE, PP_ALIGN.CENTER)

takeaways = [
    ("🏆", "最低成本 (VM 隔离)", "m7g.metal + kata-fc: $21/pod/月 (90 pods)", "裸金属消除嵌套虚拟化开销, Firecracker 轻量 microVM, Graviton 价格优势", ACCENT_GREEN),
    ("💰", "最低成本 (用户态隔离)", "m7g.2x + gVisor: $17/pod/月 (14 pods)", "零 overhead, 100% 调度器理论达成, Graviton 比 Intel 便宜 23%", ACCENT_CYAN),
    ("⚠️", "kata-qemu 最贵", "m8i.2x + kata-qemu: $44/pod/月 (仅 7 pods)", "QEMU VMExit 开销大, 嵌套虚拟化下仅达 58% 理论密度, VM crash", ACCENT_RED),
    ("📦", "EBS PVC 是隐藏瓶颈", "m7g.metal + EBS PVC: 仅 28 pods ($68/pod/月)", "Nitro 共享槽位 31 = EBS + ENI, 扣除 root/devmapper/ENI 后仅 28 PVC 可用", ACCENT_ORANGE),
    ("🔄", "运行时选择影响 2-4x 成本", "同一实例, 不同运行时: $22 vs $44/pod/月", "kata-clh 比 kata-qemu 便宜 46% (同一台 m8i.2xlarge)", ACCENT_PURPLE),
]

for i, (emoji, title, detail, note, color) in enumerate(takeaways):
    y = 1.2 + i * 1.15
    add_text_box(slide, 0.7, y, 0.5, 0.5, emoji, 22, False, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, 1.3, y, 4.0, 0.4, title, 16, True, color, PP_ALIGN.LEFT)
    add_text_box(slide, 5.5, y, 7.0, 0.4, detail, 14, True, WHITE, PP_ALIGN.LEFT)
    add_text_box(slide, 5.5, y + 0.4, 7.0, 0.4, note, 10, False, DIM_GRAY, PP_ALIGN.LEFT)

add_text_box(slide, 0.5, 7.0, 12, 0.3,
    "Data source: kata-benchmark-v2 stress tests, April 2026 · Pricing: AWS On-Demand us-west-2",
    9, False, DIM_GRAY, PP_ALIGN.CENTER)

# ── Save ──
output_path = os.path.expanduser("/home/ec2-user/kata-benchmark-v2/results/sandbox-runtime-cost-analysis.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
