#!/usr/bin/env python3
"""Generate Kata Benchmark PPT report with matplotlib charts."""

import os
import csv
import io
from collections import defaultdict

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
import matplotlib.font_manager as fm
import numpy as np

# Set font for matplotlib with CJK fallback
plt.rcParams['font.family'] = ['DejaVu Sans', 'Droid Sans Fallback']
plt.rcParams['axes.unicode_minus'] = False
fm.fontManager.addfont('/usr/share/fonts/google-droid-sans-fonts/DroidSansFallbackFull.ttf')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_BLUE = RGBColor(0x1C, 0x28, 0x33)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
MID_GRAY = RGBColor(0x85, 0x92, 0x9E)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)

# matplotlib colors
C_DARK_BLUE = '#1C2833'
C_ORANGE = '#FF9800'
C_WHITE = '#FFFFFF'
C_RUNC = '#2ECC71'
C_KATA_QEMU = '#3498DB'
C_KATA_CLH = '#E74C3C'
C_LIGHT_BG = '#F8F9FA'

RESULTS_DIR = '/home/ec2-user/kata-benchmark-v2/results'
OUTPUT_PATH = '/home/ec2-user/kata-benchmark-v2/kata-benchmark-report.pptx'
CHART_DIR = '/tmp/kata_charts'
os.makedirs(CHART_DIR, exist_ok=True)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def set_slide_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.font.name = 'Microsoft YaHei'
    return p


def add_shape_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_orange_bar(slide, top=Inches(1.15)):
    add_shape_rect(slide, Inches(0.8), top, Inches(1.5), Inches(0.06), ORANGE)


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                title, font_size=32, bold=True, color=WHITE)
    add_orange_bar(slide)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
                    subtitle, font_size=16, color=MID_GRAY)


def save_chart(fig, name):
    path = os.path.join(CHART_DIR, f'{name}.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor=fig.get_facecolor(),
                edgecolor='none', transparent=False)
    plt.close(fig)
    return path


def style_chart(ax, fig, title=''):
    fig.patch.set_facecolor(C_LIGHT_BG)
    ax.set_facecolor(C_LIGHT_BG)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#BDC3C7')
    ax.spines['bottom'].set_color('#BDC3C7')
    ax.tick_params(colors='#2C3E50', labelsize=10)
    if title:
        ax.set_title(title, fontsize=14, fontweight='bold', color=C_DARK_BLUE, pad=12)


def read_csv(filename):
    path = os.path.join(RESULTS_DIR, filename)
    with open(path, 'r') as f:
        return list(csv.DictReader(f))


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Cover
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK_BLUE)

# Top accent line
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ORANGE)

# Title
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
            'Kata Containers', font_size=52, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0),
            '\u5d4c\u5957\u865a\u62df\u5316\u6027\u80fd\u57fa\u51c6\u6d4b\u8bd5',
            font_size=40, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.8),
            'EKS \u96c6\u7fa4\u73af\u5883\u4e0b runc vs kata-qemu vs kata-clh \u5bf9\u6bd4\u8bc4\u4f30',
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Orange divider
add_shape_rect(slide, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.04), ORANGE)

# Date and info
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
            '2026-04-03  |  EKS 1.34  |  Kata Containers + OpenClaw',
            font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom accent
add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), ORANGE)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: Test Objectives
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, '\u6d4b\u8bd5\u76ee\u6807\u4e0e\u80cc\u666f')

# Left column: Why
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                 '\u4e3a\u4ec0\u4e48\u8981\u505a\u8fd9\u4e2a\u6d4b\u8bd5\uff1f', font_size=22, bold=True, color=ORANGE)
bullets = [
    '\u751f\u4ea7\u90e8\u7f72\u8bc4\u4f30 \u2014 Kata Containers \u80fd\u5426\u6ee1\u8db3\u751f\u4ea7\u73af\u5883\u6027\u80fd\u8981\u6c42',
    '\u5ba2\u6237\u5b89\u5168\u9700\u6c42 \u2014 \u5185\u6838\u7ea7\u522b\u9694\u79bb\uff0c\u9632\u6b62\u5bb9\u5668\u9003\u9038',
    '\u8d85\u5356\u6210\u672c\u4f18\u5316 \u2014 \u7a7a\u95f2\u5de5\u4f5c\u8d1f\u8f7d\u53ef\u5426\u8d85\u5356\u4ee5\u964d\u4f4e\u6210\u672c',
    '\u8fd0\u884c\u65f6\u9009\u578b \u2014 kata-qemu vs kata-clh \u54ea\u4e2a\u66f4\u9002\u5408'
]
for b in bullets:
    add_paragraph(tf, f'  \u25b8  {b}', font_size=15, color=WHITE)

# Right column: Dimensions
tf2 = add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5),
                  '\u8bc4\u4f30\u7ef4\u5ea6', font_size=22, bold=True, color=ORANGE)
dims = [
    ('\u542f\u52a8\u65f6\u95f4', '\u51b7\u542f\u52a8 / \u9971\u548c\u8282\u70b9 / \u96c6\u7fa4\u6ee1\u8f7d'),
    ('\u8d44\u6e90\u5f00\u9500', 'CPU / \u5185\u5b58 / VM \u989d\u5916\u5f00\u9500'),
    ('\u7a33\u5b9a\u6027', '\u8d85\u5356\u573a\u666f\u4e0b 2 \u5c0f\u65f6\u8fde\u7eed\u76d1\u63a7'),
    ('\u529f\u80fd\u9a8c\u8bc1', 'Gateway \u5065\u5eb7 / \u5185\u6838\u9694\u79bb\u786e\u8ba4'),
]
for title, desc in dims:
    add_paragraph(tf2, '', font_size=6, color=WHITE)
    p = add_paragraph(tf2, f'  \u25c6  {title}', font_size=16, bold=True, color=WHITE)
    add_paragraph(tf2, f'      {desc}', font_size=13, color=MID_GRAY)

# Bottom box with test scope
add_shape_rect(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), RGBColor(0x1A, 0x3A, 0x4A))
add_textbox(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.7),
            '\u6d4b\u8bd5\u8303\u56f4\uff1a5 \u4e2a\u6d4b\u8bd5\u573a\u666f  |  3 \u79cd\u8fd0\u884c\u65f6  |  9+2 \u4e2a\u8282\u70b9  |  120+ Pods  |  2 \u5c0f\u65f6\u8d85\u5356\u76d1\u63a7',
            font_size=15, color=ORANGE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Architecture Diagram
# ═══════════════════════════════════════════════════════════════
fig, ax = plt.subplots(1, 1, figsize=(12, 5.5))
fig.patch.set_facecolor(C_LIGHT_BG)
ax.set_facecolor(C_LIGHT_BG)
ax.set_xlim(0, 12)
ax.set_ylim(0, 6)
ax.axis('off')

def draw_box(ax, x, y, w, h, label, color, text_color='white', fontsize=11, alpha=1.0):
    rect = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.1",
                          facecolor=color, edgecolor='#34495E', linewidth=1.5, alpha=alpha)
    ax.add_patch(rect)
    ax.text(x + w/2, y + h/2, label, ha='center', va='center',
            fontsize=fontsize, fontweight='bold', color=text_color)

# EKS Cluster outer box
outer = FancyBboxPatch((0.3, 0.3), 11.4, 5.2, boxstyle="round,pad=0.15",
                       facecolor='none', edgecolor=C_ORANGE, linewidth=2.5, linestyle='--')
ax.add_patch(outer)
ax.text(6, 5.8, 'EKS Cluster (test-s4, us-west-2)', ha='center', va='center',
        fontsize=14, fontweight='bold', color=C_DARK_BLUE)

# runc path (left)
draw_box(ax, 0.8, 3.5, 4.8, 1.3, '', '#27AE60', alpha=0.15)
ax.text(3.2, 4.5, 'runc \u8def\u5f84', ha='center', fontsize=12, fontweight='bold', color='#27AE60')
draw_box(ax, 1.0, 3.6, 2.0, 0.9, 'containerd\n2.1.5', '#27AE60', fontsize=10)
ax.annotate('', xy=(3.2, 4.05), xytext=(3.05, 4.05),
            arrowprops=dict(arrowstyle='->', color='#27AE60', lw=2))
draw_box(ax, 3.3, 3.6, 2.0, 0.9, 'OpenClaw\nPod (runc)', '#27AE60', fontsize=10)

# kata path (right)
draw_box(ax, 6.2, 2.0, 5.3, 3.0, '', '#3498DB', alpha=0.1)
ax.text(8.85, 4.7, 'Kata \u8def\u5f84', ha='center', fontsize=12, fontweight='bold', color='#3498DB')
draw_box(ax, 6.5, 3.6, 2.0, 0.9, 'containerd\n2.1.5', '#3498DB', fontsize=10)
ax.annotate('', xy=(8.7, 4.05), xytext=(8.55, 4.05),
            arrowprops=dict(arrowstyle='->', color='#3498DB', lw=2))
draw_box(ax, 8.8, 3.6, 2.4, 0.9, 'Kata VM\n(QEMU/CLH)', '#E67E22', fontsize=10)
ax.annotate('', xy=(10.0, 3.55), xytext=(10.0, 3.0),
            arrowprops=dict(arrowstyle='->', color='#E67E22', lw=2))
draw_box(ax, 8.8, 2.1, 2.4, 0.85, 'OpenClaw\nPod (VM)', '#E74C3C', fontsize=10)

# Nodes at bottom
draw_box(ax, 0.8, 0.5, 3.3, 0.9, 'm8i.4xlarge x9\n16 vCPU / 64GB', '#2C3E50', fontsize=9)
draw_box(ax, 4.5, 0.5, 3.3, 0.9, 'm8i.4xlarge x1\n(\u65e0\u6c61\u70b9, \u5bf9\u6bd4\u7528)', '#34495E', fontsize=9)
draw_box(ax, 8.2, 0.5, 3.3, 0.9, 'r8i.2xlarge x1\n8 vCPU / 64GB (\u8d85\u5356)', '#8E44AD', fontsize=9)

ax.text(6, 1.7, '\u2191 \u8282\u70b9\u5206\u914d', ha='center', fontsize=10, color='#7F8C8D')

chart_path = save_chart(fig, 'architecture')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, '\u6d4b\u8bd5\u67b6\u6784\u56fe')
slide.shapes.add_picture(chart_path, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.4))


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Test Environment
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, '\u6d4b\u8bd5\u73af\u5883')

# Cluster info - left
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.5),
                 '\u96c6\u7fa4\u4fe1\u606f', font_size=20, bold=True, color=ORANGE)
cluster_items = [
    ('EKS \u96c6\u7fa4', 'test-s4, us-west-2'),
    ('K8s \u7248\u672c', 'v1.34.4-eks'),
    ('containerd', '2.1.5  \u26a0\ufe0f 2.2.x \u6709 Kata \u517c\u5bb9 bug'),
    ('Host Kernel', '6.12.68-92.122.amzn2023'),
    ('Kata VM Kernel', '6.18.12'),
    ('Operator', 'OpenClaw v0.22.2'),
]
for k, v in cluster_items:
    add_paragraph(tf, f'  {k}:  {v}', font_size=13, color=WHITE)

# Node config - right
tf2 = add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
                  '\u8282\u70b9\u914d\u7f6e', font_size=20, bold=True, color=ORANGE)
node_items = [
    ('m8i.4xlarge x 9', '\u2502  16 vCPU / 64GB  \u2502  \u57fa\u51c6\u6d4b\u8bd5\u8282\u70b9'),
    ('m8i.4xlarge x 1', '\u2502  16 vCPU / 64GB  \u2502  \u65e0\u6c61\u70b9\u5bf9\u6bd4\u8282\u70b9'),
    ('r8i.2xlarge x 1', '\u2502   8 vCPU / 64GB  \u2502  \u8d85\u5356\u7a33\u5b9a\u6027\u6d4b\u8bd5'),
]
for item, desc in node_items:
    add_paragraph(tf2, f'  \u25b8 {item}', font_size=14, bold=True, color=WHITE)
    add_paragraph(tf2, f'     {desc}', font_size=12, color=MID_GRAY)

# Runtime table at bottom
tf3 = add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
                  '\u8fd0\u884c\u65f6\u914d\u7f6e', font_size=20, bold=True, color=ORANGE)
rt_items = [
    'runc \u2014 \u6807\u51c6\u5bb9\u5668\u8fd0\u884c\u65f6\uff08\u57fa\u51c6\u5bf9\u7167\u7ec4\uff09',
    'kata-qemu \u2014 QEMU \u5fae\u865a\u62df\u673a\uff08\u5185\u6838\u9694\u79bb\uff09',
    'kata-clh \u2014 Cloud Hypervisor \u5fae\u865a\u62df\u673a\uff08\u5185\u6838\u9694\u79bb\uff09',
]
for item in rt_items:
    add_paragraph(tf3, f'    \u25c6  {item}', font_size=14, color=WHITE)

# Key note
add_shape_rect(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), RGBColor(0x1A, 0x3A, 0x4A))
add_textbox(slide, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.5),
            '\u5d4c\u5957\u865a\u62df\u5316\u5df2\u542f\u7528 (Intel VMX)  |  \u6bcf\u4e2a Pod \u5305\u542b 3 \u4e2a Init + 3 \u4e2a\u8fd0\u884c\u5bb9\u5668  |  gp3 EBS 10Gi',
            font_size=13, color=ORANGE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Test 1 - Cold Boot Time
# ═══════════════════════════════════════════════════════════════
data1 = read_csv('v2-test1-boot-time.csv')

runtimes = ['runc', 'kata-qemu', 'kata-clh']
runtime_labels = ['runc', 'kata-qemu', 'kata-clh']
colors = [C_RUNC, C_KATA_QEMU, C_KATA_CLH]

fig, ax = plt.subplots(figsize=(11, 5))
style_chart(ax, fig)

x = np.arange(5)
width = 0.25

for i, rt in enumerate(runtimes):
    vals = [float(r['boot_time_sec']) for r in data1 if r['runtime'] == rt]
    bars = ax.bar(x + i*width, vals, width, label=runtime_labels[i],
                  color=colors[i], edgecolor='white', linewidth=0.5)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                f'{val:.1f}s', ha='center', va='bottom', fontsize=8, fontweight='bold',
                color=colors[i])

ax.set_xlabel('\u8fed\u4ee3\u8f6e\u6b21', fontsize=12, color=C_DARK_BLUE)
ax.set_ylabel('\u542f\u52a8\u65f6\u95f4 (\u79d2)', fontsize=12, color=C_DARK_BLUE)
ax.set_xticks(x + width)
ax.set_xticklabels([f'\u7b2c {i+1} \u8f6e' for i in range(5)])
ax.legend(fontsize=11, loc='upper right')
ax.set_ylim(0, 135)
# Add average lines
for i, rt in enumerate(runtimes):
    vals = [float(r['boot_time_sec']) for r in data1 if r['runtime'] == rt]
    avg = np.mean(vals)
    ax.axhline(y=avg, color=colors[i], linestyle='--', alpha=0.4, linewidth=1)

chart_path = save_chart(fig, 'test1_boot')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, 'Test 1 \u2014 \u5355 Pod \u51b7\u542f\u52a8\u65f6\u95f4', '\u6bcf\u79cd\u8fd0\u884c\u65f6 5 \u8f6e\u6d4b\u8bd5\uff0c\u4e0d\u505a\u955c\u50cf\u9884\u62c9')
slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.8), Inches(8.5), Inches(4.2))

# Key findings on right
tf = add_textbox(slide, Inches(9.2), Inches(2.0), Inches(3.8), Inches(4.5),
                 '\u5173\u952e\u53d1\u73b0', font_size=18, bold=True, color=ORANGE)
findings = [
    '\u9996\u6b21\u51b7\u542f\u52a8\u542b VM \u955c\u50cf\u62c9\u53d6',
    'kata-qemu: ~120s',
    'kata-clh: ~107s',
    '',
    '\u70ed\u542f\u52a8\u5e73\u5747:',
    'runc: 51.75s',
    'kata-qemu: 70.45s (+37%)',
    'kata-clh: 72.04s (+39%)',
    '',
    '\u5185\u6838\u9694\u79bb\u5df2\u786e\u8ba4:',
    'Kata VM: 6.18.12',
    'Host: 6.12.68',
]
for f in findings:
    if f == '':
        add_paragraph(tf, '', font_size=6, color=WHITE)
    elif f.startswith('kata') or f.startswith('runc') or f.startswith('Host') or f.startswith('Kata VM'):
        add_paragraph(tf, f'  {f}', font_size=12, color=LIGHT_GRAY)
    else:
        add_paragraph(tf, f'  \u25b8 {f}', font_size=13, color=WHITE, bold=True)

# Bottom highlight
add_shape_rect(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.7), RGBColor(0x1A, 0x3A, 0x4A))
add_textbox(slide, Inches(0.7), Inches(6.35), Inches(12), Inches(0.6),
            'Kata \u70ed\u542f\u52a8\u5f00\u9500 +18~20\u79d2\uff08~37%\uff09\u2014 \u4e3b\u8981\u6765\u81ea VM \u521b\u5efa\u548c\u5185\u6838\u542f\u52a8\uff0c\u800c\u975e\u5e94\u7528\u5c42\u5dee\u5f02',
            font_size=15, color=ORANGE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: Test 2 & 3
# ═══════════════════════════════════════════════════════════════
data2 = read_csv('v2-test2-saturated-boot-time.csv')
data3 = read_csv('v2-test3-multi-node-boot-time.csv')

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4.5))
style_chart(ax1, fig, 'Test 2: \u9971\u548c\u8282\u70b9\u542f\u52a8 (15 pods \u5df2\u5b58\u5728)')
style_chart(ax2, fig, 'Test 3: \u96c6\u7fa4\u6ee1\u8f7d\u542f\u52a8 (120 pods \u5df2\u5b58\u5728)')

for ax, data, title_num in [(ax1, data2, 2), (ax2, data3, 3)]:
    x = np.arange(3)
    for i, rt in enumerate(runtimes):
        vals = [float(r['boot_time_sec']) for r in data if r['runtime'] == rt]
        bars = ax.bar(x + i*width, vals, width, label=runtime_labels[i],
                      color=colors[i], edgecolor='white', linewidth=0.5)
        for bar, val in zip(bars, vals):
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
                    f'{val:.0f}s', ha='center', va='bottom', fontsize=8, fontweight='bold',
                    color=colors[i])
    ax.set_xlabel('\u8fed\u4ee3\u8f6e\u6b21', fontsize=10, color=C_DARK_BLUE)
    ax.set_ylabel('\u542f\u52a8\u65f6\u95f4 (\u79d2)', fontsize=10, color=C_DARK_BLUE)
    ax.set_xticks(x + width)
    ax.set_xticklabels([f'\u7b2c {i+1} \u8f6e' for i in range(3)])
    ax.set_ylim(0, 120)
    ax.legend(fontsize=9)

fig.tight_layout(pad=2)
chart_path = save_chart(fig, 'test2_3')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, 'Test 2 & 3 \u2014 \u8282\u70b9\u9971\u548c & \u96c6\u7fa4\u6ee1\u8f7d')
slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.0))

# Summary table below
tf = add_textbox(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(1.2),
                 'Test 2 \u5e73\u5747:  runc 54.1s  |  kata-qemu 85.6s  |  kata-clh 89.7s',
                 font_size=13, color=WHITE)
add_paragraph(tf, 'Test 3 \u5e73\u5747:  runc 51.3s  |  kata-qemu 68.4s  |  kata-clh 92.3s',
              font_size=13, color=WHITE)
tf2 = add_textbox(slide, Inches(7.0), Inches(6.0), Inches(5.5), Inches(1.2),
                  '\u25b8 runc \u5728\u9971\u548c\u8282\u70b9\u4e0b\u4ec5\u589e\u52a0 ~2s\uff0c\u975e\u5e38\u7a33\u5b9a', font_size=13, color=WHITE, bold=True)
add_paragraph(tf2, '\u25b8 Kata \u5728\u9971\u548c\u8282\u70b9\u4e0b\u6ce2\u52a8\u8f83\u5927 (65-104s)', font_size=13, color=WHITE, bold=True)
add_paragraph(tf2, '\u25b8 kata-clh \u5728\u65b0\u8282\u70b9\u6709\u51b7\u542f\u52a8\u60e9\u7f5a (~105s)', font_size=13, color=WHITE, bold=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: Test 4 - Runtime Comparison
# ═══════════════════════════════════════════════════════════════
data4 = read_csv('v2-test4-runtime-comparison.csv')

fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 4.5))
style_chart(ax1, fig, '\u542f\u52a8\u65f6\u95f4\u5bf9\u6bd4 (\u540c\u4e00\u8282\u70b9)')
style_chart(ax2, fig, '\u7a33\u6001\u8d44\u6e90\u4f7f\u7528')

# Boot time
boot_times = [float(r['boot_time_sec']) for r in data4]
bars1 = ax1.bar(runtime_labels, boot_times, color=colors, edgecolor='white', linewidth=1, width=0.5)
for bar, val in zip(bars1, boot_times):
    ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1,
             f'{val:.1f}s', ha='center', va='bottom', fontsize=12, fontweight='bold',
             color=C_DARK_BLUE)
ax1.set_ylabel('\u542f\u52a8\u65f6\u95f4 (\u79d2)', fontsize=11, color=C_DARK_BLUE)
ax1.set_ylim(0, 85)

# Resource usage - parse CPU (remove 'm') and Memory (remove 'Mi')
def parse_cpu(s):
    return int(s.replace('m', ''))

def parse_mem(s):
    return int(s.replace('Mi', ''))

cpus = [parse_cpu(r['cpu_usage']) for r in data4]
mems = [parse_mem(r['memory_usage']) for r in data4]

x = np.arange(3)
w = 0.35
bars_cpu = ax2.bar(x - w/2, cpus, w, label='CPU (millicores)', color=[C_RUNC, C_KATA_QEMU, C_KATA_CLH], alpha=0.7)
ax2_twin = ax2.twinx()
bars_mem = ax2_twin.bar(x + w/2, mems, w, label='Memory (Mi)', color=[C_RUNC, C_KATA_QEMU, C_KATA_CLH], alpha=0.4, hatch='//')
ax2.set_ylabel('CPU (millicores)', fontsize=10, color=C_DARK_BLUE)
ax2_twin.set_ylabel('Memory (Mi)', fontsize=10, color=C_DARK_BLUE)
ax2.set_xticks(x)
ax2.set_xticklabels(runtime_labels)
ax2.set_ylim(0, 1200)
ax2_twin.set_ylim(0, 1200)

# Legend
from matplotlib.patches import Patch
legend_elements = [Patch(facecolor='gray', alpha=0.7, label='CPU (m)'),
                   Patch(facecolor='gray', alpha=0.4, hatch='//', label='Memory (Mi)')]
ax2.legend(handles=legend_elements, fontsize=9, loc='upper right')

for bar, val in zip(bars_cpu, cpus):
    ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 15,
             f'{val}m', ha='center', fontsize=9, fontweight='bold', color=C_DARK_BLUE)
for bar, val in zip(bars_mem, mems):
    ax2_twin.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 15,
                  f'{val}Mi', ha='center', fontsize=9, fontweight='bold', color=C_DARK_BLUE)

fig.tight_layout(pad=2)
chart_path = save_chart(fig, 'test4')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, 'Test 4 \u2014 \u8fd0\u884c\u65f6\u5bf9\u6bd4', '\u540c\u4e00\u8282\u70b9 (ip-172-31-29-155) \u4e0a\u4e09\u79cd\u8fd0\u884c\u65f6\u5bf9\u6bd4')
slide.shapes.add_picture(chart_path, Inches(0.5), Inches(2.0), Inches(12.3), Inches(4.0))

# Key insight
add_shape_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.9), RGBColor(0x1A, 0x3A, 0x4A))
tf = add_textbox(slide, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.8),
                 '\u2714 \u4e09\u79cd\u8fd0\u884c\u65f6 Gateway \u5747\u8fd4\u56de HTTP 200    \u2714 \u7a33\u6001\u8d44\u6e90\u4f7f\u7528\u51e0\u4e4e\u4e00\u6837 (~1-2m CPU, ~400Mi RAM)',
                 font_size=15, color=ORANGE, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, '\u2714 Kata VM \u5f00\u9500\u4ec5\u5728\u542f\u52a8\u9636\u6bb5\uff0c\u7a33\u6001\u65f6\u53ef\u5ffd\u7565\u4e0d\u8ba1',
              font_size=15, color=ORANGE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Test 5 - Oversell Stability (Key Metrics)
# ═══════════════════════════════════════════════════════════════
data5 = read_csv('v2-test5-oversell-stability.csv')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, 'Test 5 \u2014 \u8d85\u5356\u7a33\u5b9a\u6027 \u2b50', 'r8i.2xlarge (8 vCPU) \u4e0a\u8fd0\u884c 16 \u4e2a kata-qemu VMs\uff0c200% CPU \u8d85\u5356')

# Big metric cards
metrics = [
    ('16', 'Pods \u7a33\u5b9a\u8fd0\u884c', ORANGE),
    ('0', 'OOM \u4e8b\u4ef6', RGBColor(0x2E, 0xCC, 0x71)),
    ('2h+', '\u6301\u7eed\u76d1\u63a7', RGBColor(0x34, 0x98, 0xDB)),
    ('24', '\u91c7\u6837\u6b21\u6570', MID_GRAY),
]

for idx, (val, label, color) in enumerate(metrics):
    left = Inches(0.8 + idx * 3.1)
    add_shape_rect(slide, left, Inches(2.0), Inches(2.7), Inches(1.6), RGBColor(0x1A, 0x3A, 0x4A))
    add_textbox(slide, left, Inches(2.1), Inches(2.7), Inches(1.0),
                val, font_size=48, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, Inches(3.0), Inches(2.7), Inches(0.5),
                label, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# Restart pie chart
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(11, 3.5))
fig.patch.set_facecolor(C_LIGHT_BG)

# Pie: restart distribution
restart_counts = {0: 6, 1: 7, 2: 3}
labels = ['0 \u6b21\u91cd\u542f (6 pods)', '1 \u6b21\u91cd\u542f (7 pods)', '2 \u6b21\u91cd\u542f (3 pods)']
pie_colors = ['#2ECC71', '#F39C12', '#E74C3C']
ax1.pie(restart_counts.values(), labels=labels, colors=pie_colors,
        autopct='%1.0f%%', startangle=90, textprops={'fontsize': 10, 'color': C_DARK_BLUE})
ax1.set_title('\u91cd\u542f\u5206\u5e03', fontsize=13, fontweight='bold', color=C_DARK_BLUE)

# Time series: node CPU and memory over checks
checks = sorted(set(int(r['check_num']) for r in data5))
node_cpus = []
node_mems = []
for c in checks:
    rows = [r for r in data5 if int(r['check_num']) == c]
    if rows:
        cpu_str = rows[0]['node_cpu']
        mem_str = rows[0]['node_memory']
        node_cpus.append(int(cpu_str.replace('m', '')))
        node_mems.append(int(mem_str.replace('Mi', '')))

ax2.set_facecolor(C_LIGHT_BG)
ax2.plot(checks[:len(node_cpus)], node_cpus, color=C_ORANGE, linewidth=2, marker='o',
         markersize=4, label='\u8282\u70b9 CPU (m)')
ax2.set_ylabel('\u8282\u70b9 CPU (millicores)', fontsize=10, color=C_DARK_BLUE)
ax2.set_xlabel('\u91c7\u6837\u68c0\u67e5\u70b9', fontsize=10, color=C_DARK_BLUE)
ax2.set_ylim(0, 9000)
ax2_t = ax2.twinx()
ax2_t.plot(checks[:len(node_mems)], node_mems, color=C_KATA_QEMU, linewidth=2, marker='s',
           markersize=4, label='\u8282\u70b9 Memory (Mi)')
ax2_t.set_ylabel('\u8282\u70b9 Memory (Mi)', fontsize=10, color=C_DARK_BLUE)
ax2_t.set_ylim(15000, 25000)
ax2.spines['top'].set_visible(False)
ax2_t.spines['top'].set_visible(False)
ax2.set_title('\u8282\u70b9\u8d44\u6e90\u8d8b\u52bf (2\u5c0f\u65f6)', fontsize=13, fontweight='bold', color=C_DARK_BLUE)

lines1, labels1 = ax2.get_legend_handles_labels()
lines2, labels2 = ax2_t.get_legend_handles_labels()
ax2.legend(lines1 + lines2, labels1 + labels2, fontsize=9, loc='center right')

fig.tight_layout(pad=2)
chart_path = save_chart(fig, 'test5')
slide.shapes.add_picture(chart_path, Inches(0.5), Inches(3.8), Inches(12.3), Inches(3.5))


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Key Findings & Recommendations
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, '\u5173\u952e\u53d1\u73b0\u4e0e\u5efa\u8bae')

# Findings - left
tf = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.8), Inches(5.0),
                 '\u5173\u952e\u53d1\u73b0', font_size=22, bold=True, color=ORANGE)

findings = [
    ('\u542f\u52a8\u5f00\u9500 +37%', 'Kata \u70ed\u542f\u52a8\u6bd4 runc \u6162 ~18-20 \u79d2\uff0c\u4e3b\u8981\u6765\u81ea VM \u521b\u5efa'),
    ('\u7a33\u6001\u8d44\u6e90\u51e0\u4e4e\u4e00\u6837', '\u7a7a\u95f2\u65f6 CPU/\u5185\u5b58\u5f00\u9500\u4e0e runc \u65e0\u660e\u663e\u5dee\u5f02'),
    ('\u8d85\u5356\u53ef\u884c', '200% CPU \u8d85\u5356\u4e0b 16 VMs \u7a33\u5b9a\u8fd0\u884c 2h+\uff0c0 OOM'),
    ('\u5185\u6838\u9694\u79bb\u786e\u8ba4', 'Kata VM \u8fd0\u884c\u72ec\u7acb\u5185\u6838 6.18.12\uff0c\u4e0e Host \u5b8c\u5168\u9694\u79bb'),
    ('kata-qemu \u2248 kata-clh', '\u6027\u80fd\u76f8\u8fd1\uff0ckata-qemu \u66f4\u7a33\u5b9a\uff0c\u51b7\u542f\u52a8\u6ce2\u52a8\u66f4\u5c0f'),
]
for title, desc in findings:
    add_paragraph(tf, '', font_size=4, color=WHITE)
    add_paragraph(tf, f'  \u25b8 {title}', font_size=16, bold=True, color=WHITE)
    add_paragraph(tf, f'     {desc}', font_size=13, color=MID_GRAY)

# Recommendations - right
tf2 = add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.8), Inches(5.0),
                  '\u5efa\u8bae', font_size=22, bold=True, color=ORANGE)

recs = [
    ('\u751f\u4ea7\u73af\u5883\u63a8\u8350 kata-qemu', '\u6bd4 kata-clh \u66f4\u7a33\u5b9a\uff0c\u51b7\u542f\u52a8\u8868\u73b0\u66f4\u53ef\u9884\u6d4b'),
    ('\u5206\u6279\u542f\u52a8', '\u907f\u514d\u542f\u52a8\u98ce\u66b4\u5bfc\u81f4 startup probe \u5931\u8d25'),
    ('\u8d85\u5356\u7a7a\u95f2 workload', '200% CPU \u8d85\u5356\u5bf9\u7a7a\u95f2\u5de5\u4f5c\u8d1f\u8f7d\u5b89\u5168\u53ef\u884c'),
    ('\u9884\u7559\u542f\u52a8 buffer', 'Kata \u542f\u52a8\u9700 70-90s\uff08vs runc 50s\uff09'),
    ('containerd \u7248\u672c\u9501\u5b9a', '\u5fc5\u987b\u4f7f\u7528 2.1.x\uff0c2.2.x \u6709 Kata \u517c\u5bb9 bug'),
]
for title, desc in recs:
    add_paragraph(tf2, '', font_size=4, color=WHITE)
    add_paragraph(tf2, f'  \u25c6 {title}', font_size=16, bold=True, color=WHITE)
    add_paragraph(tf2, f'     {desc}', font_size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Summary & Next Steps
# ═══════════════════════════════════════════════════════════════

# Summary comparison chart
fig, ax = plt.subplots(figsize=(10, 4))
style_chart(ax, fig, '\u5404\u573a\u666f\u542f\u52a8\u65f6\u95f4\u6c47\u603b\u5bf9\u6bd4')

scenarios = ['\u51b7\u542f\u52a8\n(\u9996\u6b21)', '\u70ed\u542f\u52a8\n(\u5e73\u5747)', '\u9971\u548c\u8282\u70b9', '\u96c6\u7fa4\u6ee1\u8f7d']
runc_vals = [49.22, 51.75, 54.09, 51.26]
kata_qemu_vals = [119.82, 70.45, 85.57, 68.43]
kata_clh_vals = [107.10, 72.04, 89.74, 92.28]

x = np.arange(len(scenarios))
w = 0.25

for i, (vals, color, label) in enumerate([
    (runc_vals, C_RUNC, 'runc'),
    (kata_qemu_vals, C_KATA_QEMU, 'kata-qemu'),
    (kata_clh_vals, C_KATA_CLH, 'kata-clh')
]):
    bars = ax.bar(x + i*w, vals, w, label=label, color=color, edgecolor='white')
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 1.5,
                f'{val:.0f}s', ha='center', fontsize=9, fontweight='bold', color=color)

ax.set_ylabel('\u542f\u52a8\u65f6\u95f4 (\u79d2)', fontsize=11, color=C_DARK_BLUE)
ax.set_xticks(x + w)
ax.set_xticklabels(scenarios, fontsize=10)
ax.set_ylim(0, 140)
ax.legend(fontsize=11, loc='upper right')

chart_path = save_chart(fig, 'summary')

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BLUE)
add_slide_title(slide, '\u603b\u7ed3\u4e0e\u4e0b\u4e00\u6b65')

slide.shapes.add_picture(chart_path, Inches(0.5), Inches(1.8), Inches(7.5), Inches(3.2))

# Right side: production recommendation
tf = add_textbox(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(3.5),
                 '\u751f\u4ea7\u90e8\u7f72\u5efa\u8bae', font_size=22, bold=True, color=ORANGE)
add_paragraph(tf, '', font_size=6, color=WHITE)
add_paragraph(tf, '  \u63a8\u8350\u8fd0\u884c\u65f6: kata-qemu', font_size=18, bold=True, color=RGBColor(0x2E, 0xCC, 0x71))
add_paragraph(tf, '', font_size=6, color=WHITE)
add_paragraph(tf, '  \u2714 \u6027\u80fd\u6700\u7a33\u5b9a\uff0c\u6ce2\u52a8\u6700\u5c0f', font_size=14, color=WHITE)
add_paragraph(tf, '  \u2714 \u5185\u6838\u7ea7\u522b\u5b89\u5168\u9694\u79bb', font_size=14, color=WHITE)
add_paragraph(tf, '  \u2714 \u7a33\u6001\u8d44\u6e90\u5f00\u9500\u4e0e runc \u76f8\u5f53', font_size=14, color=WHITE)
add_paragraph(tf, '  \u2714 \u8d85\u5356\u573a\u666f\u5df2\u9a8c\u8bc1\u53ef\u884c', font_size=14, color=WHITE)

# Next steps at bottom
add_shape_rect(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.9), RGBColor(0x1A, 0x3A, 0x4A))
tf2 = add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
                  '\u4e0b\u4e00\u6b65\u8ba1\u5212', font_size=20, bold=True, color=ORANGE)
next_steps = [
    '\u9ad8\u8d1f\u8f7d\u538b\u6d4b \u2014 \u6d4b\u8bd5 Kata VM \u5728\u9ad8 CPU/\u5185\u5b58\u8d1f\u8f7d\u4e0b\u7684\u8868\u73b0\u548c\u7a33\u5b9a\u6027',
    '\u7f51\u7edc\u6027\u80fd\u6d4b\u8bd5 \u2014 Kata VM \u7f51\u7edc\u5ef6\u8fdf/\u541e\u5410\u91cf vs runc',
    '\u81ea\u52a8\u5316\u96c6\u6210 \u2014 \u5c06 Kata \u90e8\u7f72\u96c6\u6210\u5230 CI/CD \u6d41\u6c34\u7ebf',
    '\u751f\u4ea7\u73af\u5883 Pilot \u2014 \u9009\u62e9\u4f4e\u98ce\u9669\u5de5\u4f5c\u8d1f\u8f7d\u5f00\u59cb\u8bd5\u70b9',
]
for step in next_steps:
    add_paragraph(tf2, f'  \u25b8 {step}', font_size=14, color=WHITE)

# Bottom accent
add_shape_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), ORANGE)


# ── Save ──
prs.save(OUTPUT_PATH)
print(f'PPT saved to: {OUTPUT_PATH}')
